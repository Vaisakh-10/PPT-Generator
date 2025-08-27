import streamlit as st
import openai
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt

# -------------------------
# Helper: Extract layouts from template
# -------------------------
def load_template(template_file):
    return Presentation(template_file)

# -------------------------
# Helper: Call LLM to split text into slides
# -------------------------
def generate_slide_outline(api_key, text, guidance):
    openai.api_key = api_key
    prompt = f"""
    You are an assistant that converts text into PowerPoint slide outlines.
    Input text:
    {text}

    Guidance: {guidance}

    Output JSON with this format:
    [
      {{"title": "Slide 1 Title", "content": ["bullet point 1", "bullet point 2"]}},
      {{"title": "Slide 2 Title", "content": ["bullet point 1", "bullet point 2"]}}
    ]
    """
    response = openai.ChatCompletion.create(
        model="gpt-4o-mini",
        messages=[{"role":"user","content":prompt}],
        temperature=0.3
    )
    import json
    try:
        slides = json.loads(response.choices[0].message["content"])
    except Exception:
        slides = [{"title":"Error parsing response","content":[response.choices[0].message["content"]]}]
    return slides

# -------------------------
# Helper: Build new PPT from outline + template
# -------------------------
def build_presentation(slides, template_file):
    prs = Presentation(template_file)
    title_layout = prs.slide_layouts[0]  # usually Title Slide
    bullet_layout = prs.slide_layouts[1] # usually Title + Content

    # Remove existing slides (if template has many)
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    # Add slides
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(bullet_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = slide_data.get("title","")
        tf = body.text_frame
        tf.clear()
        for bullet in slide_data.get("content",[]):
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)

    # Save to temp file
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name

# -------------------------
# Streamlit App
# -------------------------
st.title("📊 Auto PPT Generator")
st.write("Turn your text/markdown into a styled PowerPoint using your own template.")

user_text = st.text_area("Paste your text here:", height=200)
guidance = st.text_input("Optional guidance (e.g. 'Investor pitch deck')", "")
api_key = st.text_input("Enter your LLM API key (not stored):", type="password")
uploaded_template = st.file_uploader("Upload a PPTX/POTX template", type=["pptx","potx"])

if st.button("Generate Presentation"):
    if not (user_text and api_key and uploaded_template):
        st.error("Please provide text, API key, and template file.")
    else:
        with st.spinner("Generating slides..."):
            slides = generate_slide_outline(api_key, user_text, guidance)
            ppt_file = build_presentation(slides, uploaded_template)

        with open(ppt_file, "rb") as f:
            st.download_button("⬇️ Download PPT", f, file_name="generated.pptx")
